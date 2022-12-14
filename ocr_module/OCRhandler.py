
from docx import Document
from pandas import read_excel
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pytesseract import pytesseract


#############pytesseract.tesseract_cmd = r"C:\Users\user1\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"


class OCRHandler(object):
    def __init__(self, path_to_tesseract: str, path_to_poppler: str):
        self.path_to_tesseract = path_to_tesseract
        pytesseract.tesseract_cmd = self.path_to_tesseract
        self.path_to_poppler = path_to_poppler

    def run_handle(self, file_path: str, file_type: str):
        if file_type == "pdf":
            txt = self.pdf(file_path=file_path)
        elif file_type == "png":
            txt = self.png(file_path=file_path)
        elif file_type == "jpeg":
            txt = self.jpeg(file_path=file_path)
        elif file_type == "pptx":
            txt = self.pptx(file_path=file_path)
        elif file_type == "jpg":
            txt = self.jpg(file_path=file_path)
        elif file_type == "docx":
            txt = self.docx(file_path=file_path)
        elif file_type == "xlsx":
            txt = self.xlsx(file_path=file_path)
        elif file_type == "gif":
            txt = self.gif(file_path=file_path)
        return txt

    def pdf(self, file_path):
        text = ""
        #images = convert_from_path(file_path, poppler_path = r"C:\Users\user1\Downloads\Release-22.04.0-0\poppler-22.04.0\Library\bin")
        images = convert_from_path(file_path, poppler_path=self.path_to_poppler)
        for image in images:
            text = text + pytesseract.image_to_string(image)
            print(text)
        return text


    def png(self, file_path):
        text = pytesseract.image_to_string(file_path)
        print(text)
        return text


    def jpeg(self, file_path):
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        img.close()
        print(text)
        return text

    def pptx(self, file_path):
       text = ""
       try:
           ppt = Presentation(file_path)
           for slide in ppt.slides:
               for shape in slide.shapes:
                   if hasattr(shape, "text"):
                       text = text + " " + shape.text
           return text
       except:
           return "pptx file/file-path not found"

    def jpg(self, file_path):
        text = pytesseract.image_to_string(file_path)
        print(text)
        return text

    def docx(self, file_path):
        text = ""
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text = text + para.text
            return text
        except:
            return "docx file/file-path not found"

    def xlsx(self, file_path):
        text = ""
        try:
            df = read_excel(file_path)
            for i in range(0, df.shape[0]): # to access row
                for j in range(0, df.shape[1]): # to access each field in row
                    text = text + " " + str(df.iloc[i][j])
            return text
        except:
            return "xlsx file/file-path not found"

    def gif(self, file_path):
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            img.close()
            return text
        except:
            return "gif file/file-path not found"

